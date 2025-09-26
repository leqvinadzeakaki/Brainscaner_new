import os
import pdfplumber
from pptx import Presentation
from flask import Flask, render_template, request, session, redirect, url_for, abort
from dotenv import load_dotenv

import google.generativeai as genai
from google_auth_oauthlib.flow import Flow
import google.oauth2.credentials
import google.auth.transport.requests
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

# --- Flask Config ---
load_dotenv()

app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "dev-secret-change-me")

app.config.update(
    SESSION_COOKIE_SECURE=True,      # მხოლოდ HTTPS-ზე გაეგზავნება session cookie
    SESSION_COOKIE_SAMESITE="None",  # OAuth cross-site redirect-ზე ქუქი არ დაიკარგოს
    PREFERRED_URL_SCHEME="https",
)
app.config['UPLOAD_FOLDER'] = 'uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# --- Gemini API Config ---
# ENV-ში უნდა გქონდეს GEMINI_API_KEY=<შენი key>
genai.configure(api_key=os.environ["GEMINI_API_KEY"])
model = genai.GenerativeModel("gemini-1.5-flash")

# --- Helpers ---
def extract_text_from_pdf(path):
    text = ""
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        print(f"[PDF Error] {e}")
    return text

def extract_text_from_pptx(path):
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"[PPTX Error] {e}")
    return text

def analyze_idea(idea_text):
    prompt = f"""
შეაფასე შემდეგი ბიზნეს იდეა დეტალურად, შემდეგი სტრუქტურის მიხედვით:

ბიზნეს იდეა:
{idea_text}

1. იდეის მოკლე რეზიუმე
2. მიზნობრივი აუდიტორია
3. მონეტიზაციის გზები
4. ანალოგიური პროდუქტები ან კონკურენტები
5. იდეის სიძლიერეები და სუსტი მხარეები
6. გრძელვადიანი მდგრადობის პროგნოზი
7. რეკომენდაცია იდეის გაუმჯობესებისთვის
"""
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"❌ შეცდომა Gemini API-სთან დაკავშირებისას: {e}"

def save_analysis_to_file(content, filename="analysis.txt"):
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content)
    return filepath

def upload_to_user_drive(filepath, filename, folder_id=None):
    if "credentials" not in session:
        return None
    creds = google.oauth2.credentials.Credentials(**session['credentials'])
    service = build('drive', 'v3', credentials=creds)

    metadata = {'name': filename}
    if folder_id:
        metadata['parents'] = [folder_id]

    media = MediaFileUpload(filepath, resumable=True)
    uploaded = service.files().create(body=metadata, media_body=media, fields='id').execute()
    file_id = uploaded.get('id')

    service.permissions().create(
        fileId=file_id,
        body={'type': 'anyone', 'role': 'reader'}
    ).execute()

    return f"https://drive.google.com/file/d/{file_id}/view?usp=sharing"

def save_and_return_link(idea_text, base_filename="idea_analysis"):
    result = analyze_idea(idea_text)
    filename = f"{base_filename}.txt"
    filepath = save_analysis_to_file(result, filename)
    drive_link = upload_to_user_drive(filepath, filename)

    if 'history' not in session:
        session['history'] = []
    session['history'].append({'filename': filename, 'drive_link': drive_link})
    session['history'] = session['history'][-30:]
    return drive_link, result

# --- Health check (Render-friendly) ---
@app.route("/healthz")
def healthz():
    return "ok", 200

# --- OAuth2 ---
@app.before_request
def require_login():
    # ლოგინამდე დაშვებული endpoint-ები (root '/' აქ არ შედის!)
    allowed = ['login', 'oauth2callback', 'static', 'healthz']
    if request.endpoint in allowed or 'credentials' in session:
        return
    return redirect(url_for('login'))

@app.route("/login")
def login():
    flow = Flow.from_client_secrets_file(
        "client_secret.json",
        scopes=["https://www.googleapis.com/auth/drive.file"],
        redirect_uri=url_for("oauth2callback", _external=True, _scheme="https"),
    )
    authorization_url, state = flow.authorization_url(
        access_type="offline",
        include_granted_scopes="true",
        prompt="consent",
    )
    session["state"] = state
    return redirect(authorization_url)

@app.route("/oauth2callback")
def oauth2callback():
    state = session.get("state")

    # თუ სესიაში state არაა, არ ვაძლევთ 400-ს — ვაბრუნებთ /login-ზე, რომ სუფთად დაიწყოს flow
    if not state:
        return redirect(url_for("login"))

    flow = Flow.from_client_secrets_file(
        "client_secret.json",
        scopes=["https://www.googleapis.com/auth/drive.file"],
        state=state,
        redirect_uri=url_for("oauth2callback", _external=True, _scheme="https"),
    )
    flow.fetch_token(authorization_response=request.url)

    credentials = flow.credentials
    session["credentials"] = {
        "token": credentials.token,
        "refresh_token": credentials.refresh_token,
        "token_uri": credentials.token_uri,
        "client_id": credentials.client_id,
        "client_secret": credentials.client_secret,
        "scopes": credentials.scopes,
    }

    session.pop("state", None)
    return redirect(url_for("index"))

@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login"))

# --- Main Route ---
@app.route("/", methods=["GET", "POST"])
def index():
    result = None
    drive_link = None
    error = None

    if request.method == "POST":
        text_idea = request.form.get("text_idea")
        file = request.files.get("file")

        idea_text = ""
        if text_idea and text_idea.strip():
            idea_text = text_idea.strip()
            source_name = "text_input"
        elif file and file.filename:
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(filepath)
            if filepath.lower().endswith(".pdf"):
                idea_text = extract_text_from_pdf(filepath)
            elif filepath.lower().endswith(".pptx"):
                idea_text = extract_text_from_pptx(filepath)
            else:
                error = "❌ მხოლოდ .pdf და .pptx ფაილებია მხარდაჭერილი."
                return render_template("index.html", result=None, drive_link=None, error=error, history=session.get('history', []))
            source_name = file.filename
        else:
            error = "გთხოვთ, შეიყვანეთ ტექსტი ან ატვირთეთ ფაილი."

        if idea_text.strip():
            drive_link, result = save_and_return_link(idea_text, base_filename=source_name)

    return render_template("index.html", result=result, drive_link=drive_link, error=error, history=session.get('history', []))

# --- Run (local dev) ---
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 10000))
    app.run(debug=True, host="0.0.0.0", port=port)
